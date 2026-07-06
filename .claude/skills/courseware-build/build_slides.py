#!/usr/bin/env python3
"""Generate the IoT Fundamental for Beginners slide deck (all-white Tertiary
house style). Content is driven entirely by course_data.py + data_labs.py so
the deck stays 100% aligned with the LP, LG and labs/.

Reuses the wsq-slides reference visual components: cover, section, content,
two_col, cards3, big_statement, tile_grid, flow_h, trainer_slide,
activity_overview, step_slide, test_slide, brk, img_slide.
"""
import os, sys, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0, HERE)
import course_data as C
from data_labs import LABS

REPO = os.path.dirname(os.path.dirname(os.path.dirname(HERE)))  # .claude/skills/courseware-build -> repo root
ASSETS = os.path.join(REPO, "courseware", "assets")

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    tsz=29 if len(title)<=48 else 23   # auto-shrink long titles so they never hit the rule
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,tsz,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    rect(s,Inches(11.0),Inches(0.72),Inches(1.55),Inches(1.0),TEAL)
    txt(s,Inches(11.0),Inches(0.84),Inches(1.55),Inches(0.5),[[("IoT",22,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(11.0),Inches(1.3),Inches(1.55),Inches(0.4),[[("IOTFLOW · n8n",8,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.35),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(4.65),Inches(12),Inches(1.6),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)],
         [(f"IoT platform: {C.PLATFORM_URL}",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,20,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Uses:  ",13,GREY,True),(services,13,GREY,False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.6),[[(text,21,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        lines=cmd.split("\n"); ch=max(0.95, 0.34*len(lines)+0.42)
        rect(s,Inches(2.55),Inches(4.35),Inches(10.1),Inches(ch),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.5),Inches(9.7),Inches(ch-0.3),
            [[(ln,12,RGBColor(0x9C,0xDC,0xFE),False)] for ln in lines],space=2)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.9),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.7),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1
def img_slide(title,img,caption="",kicker="DIAGRAM"):
    """Diagram slide: image centred and scaled to fit the content area."""
    path=os.path.join(ASSETS,img)
    if not os.path.exists(path): return content(title,[caption or img],kicker=kicker)
    s=head(slide(),title,kicker)
    iw,ih=Image.open(path).size
    maxw=Inches(11.6)
    maxh=Inches(4.45) if caption else Inches(4.85)
    scale=min(maxw/iw,maxh/ih)
    w,h=int(iw*scale),int(ih*scale)
    x=int((SW-w)/2); y=Inches(1.95)+int((maxh-h)/2)
    s.shapes.add_picture(path,x,y,width=w,height=h)
    if caption:
        txt(s,Inches(0.85),Inches(6.5),Inches(11.6),Inches(0.5),[[(caption,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Expertise","IoT, data analytics, AI/ML and workflow automation (n8n)."),
  ("Delivers","WSQ courses on IoT, AI, data science and software engineering."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with IoT, electronics or automation (if any).",
 "What you want to sense, control or automate after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
content("LMS / TMS",[
 "Access your course materials, attendance and assessment on the LMS/TMS portal.",
 "Portal: https://ai-lms-tms.tertiaryinfo.tech/",
 "Download the slides and Learner Guide for reference during the open-book assessment."],kicker="COURSE PORTAL")
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Digital Attendance (AM) · Introductions",1),
 ("Topic 1: Overview of Internet of Things (IoT)",1),
 ("Topic 2: Collect and Post Data to Cloud",1),
 ("Labs 1–2 on the IoTFlow platform",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Topic 3: Read Data and Remote Control from Cloud (Labs 3–4)",1),
 ("Topic 4: IoT Data Analytics and Visualization (Labs 5–6)",1),
 ("Course Feedback and TRAQOM Survey",1),
 ("Final Assessment (WA + PP)",1),
 ("Daily timing",0),
 ("9:30am–6:30pm · 1-hour lunch · tea breaks within",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2 & timing")
tile_grid("Skills Framework (TSC)",[
 ("TSC Title",C.TSC_TITLE),
 ("TSC Code",C.TSC_CODE),
 ("Abilities","A1 briefings on IoT uses & functions · A2 review IoT testing results · A3 integrate multiple data sources · A4 review data for business value."),
 ("Knowledge","K1 IoT concept · K2 IoT devices · K3 circuits & sensors · K4 wireless technologies · K5 data analytics · K6 cybersecurity.")],
 kicker="WSQ · SKILLS FRAMEWORK",cols=1,size=14)
tile_grid("Learning Outcomes",[
 ("LO1","Understand the uses and functions of IoT technologies."),
 ("LO2","Post sensor data to cloud for IoT review."),
 ("LO3","Control devices from cloud data sources."),
 ("LO4","Data analytics and visualization on cloud to gain business insight.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
two_col("Course Outline — Day 1",
 [(f"Topic 1  {C.TOPICS[0]['title']}",0),(f"({C.TOPICS[0]['tags']})",1)]+[(o,1) for o in C.TOPICS[0]["outline"]],
 [(f"Topic 2  {C.TOPICS[1]['title']}",0),(f"({C.TOPICS[1]['tags']})",1)]+[(o,1) for o in C.TOPICS[1]["outline"]],
 kicker="4 TOPICS · 6 LABS",lhead="Morning",rhead="Afternoon")
two_col("Course Outline — Day 2",
 [(f"Topic 3  {C.TOPICS[2]['title']}",0),(f"({C.TOPICS[2]['tags']})",1)]+[(o,1) for o in C.TOPICS[2]["outline"]],
 [(f"Topic 4  {C.TOPICS[3]['title']}",0),(f"({C.TOPICS[3]['tags']})",1)]+[(o,1) for o in C.TOPICS[3]["outline"]],
 kicker="4 TOPICS · 6 LABS",lhead="Morning",rhead="Afternoon")
content("Criteria for Funding",[
 "Minimum attendance rate of 75% based on the SSG Digital Attendance record.",
 "Complete the assessment and be assessed as 'Competent'."],kicker="WSQ FUNDING")
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ---------------- OUR IOT PLATFORM ----------------
section("OUR IOT PLATFORM",f"{C.PLATFORM} — {C.PLATFORM_URL.replace('https://','')}",
        "", "The low-code IoT platform you will use for every lab in this course")
big_statement("Connect. Visualise. Automate.",C.PLATFORM_TAGLINE,
              f"{C.PLATFORM} · {C.PLATFORM_URL}",color=TEAL)
img_slide(f"{C.PLATFORM} — Low-code IoT Platform powered by n8n","iotflow-home.png",
          f"Sign up free at {C.PLATFORM_URL} — you will use it in every lab.",kicker="PLATFORM TOUR")
flow_h("How it Works",[st for st in C.HOW_IT_WORKS],kicker=f"{C.PLATFORM_URL}/#how",color=TEAL)
tile_grid(f"{C.PLATFORM} Features (1 of 2)",C.PLATFORM_FEATURES[:4],kicker="PLATFORM TOUR",cols=2,size=14)
tile_grid(f"{C.PLATFORM} Features (2 of 2)",C.PLATFORM_FEATURES[4:],kicker="PLATFORM TOUR",cols=2,size=14)
tile_grid("Industries Served",[(i,"") for i in C.INDUSTRIES],kicker="WHERE IOTFLOW IS USED",cols=3,size=13)
cards3("Platform Tutorials → Course Labs",[
 (BLUE,"Get connected",["Tutorial: Add your first device → Lab 1","Tutorial: Send your first reading → Lab 2"]),
 (TEAL,"Read & control",["Read state via REST & MQTT → Lab 3","Tutorial: Control a device → Lab 4"]),
 (VIOLET,"Visualise & automate",["Tutorial: Build a dashboard → Lab 5","Tutorial: Automate with n8n → Lab 6"])],
 kicker=C.TUTORIALS_URL)

# ---------------- TOPICS + LABS ----------------
TOPIC_LABS={t["num"]:[a for a in LABS if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]

def labs_cards(t):
    acts=TOPIC_LABS[t["num"]]
    if not acts: return
    tile_grid(f"Hands-On Labs — Topic {t['code']}",
              [(f"Lab {a['num']} — {a['title']}",a["build"]) for a in acts],
              kicker=t["title"].upper(),cols=1,size=15,
              icons=[str(a["num"]) for a in acts])

def emit_labs(t):
    for a in TOPIC_LABS[t["num"]]:
        activity_overview(f"LAB {a['num']}",a["title"],a["desc"],a["build"],a["services"],
                          kicker=f"TOPIC {t['code']} · HANDS-ON")
        if a.get("shot"):
            img_slide(a["title"],a["shot"],
                      f"Follow along on the platform: {a['tutorial']}",
                      kicker=f"LAB {a['num']} · PLATFORM WALKTHROUGH")
        steps=a["steps"]; total=len(steps)
        for i,(instr,cmd) in enumerate(steps,1):
            step_slide(f"LAB {a['num']} · STEP-BY-STEP",a["title"],i,total,instr,cmd)
        test_slide(a["title"],a["test"],kicker=f"LAB {a['num']} · VERIFY")

# ================= TOPIC 1 =================
T=C.TOPICS[0]
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"ABILITIES & KNOWLEDGE: {T['tags']}",cols=2,size=14)
img_slide("The 4th Industrial Revolution","industry40.png",
          "Industry 4.0: cyber-physical systems — IoT, big data, cloud, AI and automation.",kicker="CONTEXT")
content("What is the Internet of Things (IoT)?",[
 "IoT refers to the network of physical devices, vehicles, buildings and other objects embedded with sensors, software and connectivity.",
 "These 'things' collect and exchange data over the internet — without any human interaction.",
 "Low-cost sensors and long-range wireless technology are instrumental to IoT; 5G drives mass adoption.",
 "By connecting the physical world to the cloud, IoT turns real-world events into data you can analyse and act on."],kicker="DEFINITION")
img_slide("Internet of Things — Everything Connected","iot-network.png",
          "Devices, people and platforms exchanging data over the internet.",kicker="WHAT IS IOT")
img_slide("IoT Components","iot-components.png",
          "A 'thing' = embedded electronics + sensors + connectivity + data.",kicker="BUILDING BLOCKS")
img_slide("How Does IoT Work?","how-iot-works.png",
          "Sensing devices stream data to an IoT platform, which integrates, analyses and shares the most valuable information.",kicker="END-TO-END")
tile_grid("What Are IoT Devices?",[
 ("Physical objects that sense","Hardware such as sensors, gadgets, appliances and machines that collect and exchange data over the internet."),
 ("Programmed for a purpose","Each device is built for a specific application and can be embedded into other IoT devices."),
 ("Connected & addressable","Integrated CPU, network adapter and firmware, with an IP address on the network."),
 ("Managed by software","Configured and controlled through an app or dashboard — e.g. your phone controls the lights at home.")],
 kicker="K2 · IOT DEVICES",cols=2,size=14)
img_slide("Sensors for IoT","sensors-overview.png",
          "Temperature, humidity, IR, proximity, gas, chemical, motion, light — sensors are the device's senses.",kicker="SENSING")
img_slide("More Sensors for IoT","sensors-more.png",
          "Pick the sensor for the physical quantity you need to measure.",kicker="SENSING")
img_slide("IoT Actuators","iot-actuator.png",
          "Actuators act on the world: relays, motors, pumps, valves, LEDs — the output side of IoT.",kicker="ACTING")
tile_grid("Sensors → Triggers → Actions",[
 ("Sense","A sensor produces a reading — temperature 31 °C, door open, tank level low."),
 ("Trigger","A rule watches the reading — 'temperature above 30' or 'device offline' fires an event."),
 ("Act","The event drives an action — notify someone, log it, call AI, or command an actuator."),
 ("Automate","On IoTFlow, triggers hand off to n8n workflows — the full loop with zero code (Topic 4).")],
 kicker="THE CORE IOT PATTERN",cols=2,size=14)
content("Wireless Technologies for IoT (Short Range)",[
 "Wi-Fi — widely used, high data rates and reliable connectivity over short-to-medium distances.",
 "Bluetooth — low-power, short-range; connects IoT devices to phones and tablets.",
 "Zigbee — low-power wireless mesh for low data rates and reliable longer-distance links.",
 "Z-Wave — optimised for low-power devices; reliable communication over longer distances."],kicker="K4 · CONNECTIVITY",size=18)
content("Wireless Technologies for IoT (Long Range)",[
 "LoRaWAN — long-range, low-power; designed for smart-city-scale applications.",
 "Sigfox — LPWAN for low data rates and long range, e.g. asset tracking and remote monitoring.",
 "NB-IoT — narrowband cellular LPWAN with reliable long-range connectivity.",
 "5G — high bandwidth and massive device density; the driver for IoT mass adoption."],kicker="K4 · CONNECTIVITY",size=18)
img_slide("Comparison of Wireless Technologies","wireless-comparison.png",
          "Trade-off: range vs bandwidth vs power consumption.",kicker="K4 · CONNECTIVITY")
tile_grid("Types of IoT",[
 ("Consumer IoT","Everyday use — home appliances, voice assistants, light fixtures."),
 ("Commercial IoT","Healthcare and transport — smart pacemakers, monitoring, connected cars."),
 ("Industrial IoT (IIoT)","Manufacturing and energy — digital control systems, smart agriculture, industrial big data."),
 ("Infrastructure IoT","Smart-city connectivity — infrastructure sensors and management systems.")],
 kicker="CLASSIFICATION",cols=2,size=14)
img_slide("IoT vs Industrial IoT (IIoT)","iot-vs-iiot.png",
          "IIoT extends IoT into industrial sectors — real-time data from sensors to equipment improves business processes.",kicker="IIOT")
img_slide("IoT Applications","iot-applications.png",
          "IoT has found its way into every conceivable field — from agriculture to aerospace.",kicker="USE CASES")
img_slide("Use Case — Healthcare","usecase-healthcare.png",
          "Track patient-assistance assets: IoT-tagged wheelchairs located in seconds.",kicker="USE CASES")
img_slide("Use Case — Smart Home","usecase-smart-home.png",
          "Smart outlets, thermostats, cameras and locks — efficiency, safety and comfort.",kicker="USE CASES")
img_slide("Use Case — Agriculture","usecase-agriculture.png",
          "Soil moisture, weather and crop sensors drive precision irrigation.",kicker="USE CASES")
img_slide("Use Case — Manufacturing","usecase-manufacturing.png",
          "Production-line monitoring enables proactive maintenance before failure.",kicker="USE CASES")
img_slide("Use Case — Logistics & Smart City","usecase-logistics.png",
          "Fleets rerouted on live conditions; cities manage traffic, utilities and environment.",kicker="USE CASES")
content(f"Recap — {T['title']}",[
 "IoT connects physical devices with sensors, software and connectivity to collect and exchange data.",
 "Devices sense with sensors, act with actuators, and triggers turn readings into events and actions.",
 "Short-range (Wi-Fi, BLE, Zigbee) and long-range (LoRaWAN, NB-IoT, 5G) wireless link devices to the cloud.",
 "IoT use cases span smart home, healthcare, agriculture, manufacturing, logistics and smart cities."],kicker="TOPIC RECAP",size=17)

brk("Lunch Break","1 hour")

# ================= TOPIC 2 =================
T=C.TOPICS[1]
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"ABILITIES & KNOWLEDGE: {T['tags']}",cols=2,size=14)
img_slide("Cloud Computing for IoT","cloud-computing.png",
          "The cloud stores, processes and serves IoT data at scale — pay-as-you-go, no servers to run.",kicker="CLOUD")
tile_grid("IoT Cloud Platforms",[
 ("IoTFlow (this course)","Our low-code platform — MQTT/HTTP, dashboards, virtual pins, n8n automation, AI."),
 ("ThingSpeak","MATLAB-centric IoT analytics platform with channels and REST/MQTT APIs."),
 ("AWS IoT / Azure IoT / Google","Hyperscaler IoT suites for enterprise fleets and big-data pipelines."),
 ("Blynk / Arduino Cloud","Maker-friendly device platforms with mobile apps and virtual pins.")],
 kicker="THE LANDSCAPE",cols=2,size=14)
content("Two Ways to Post Data to the Cloud",[
 "HTTP REST API — the device POSTs a JSON document to an endpoint; simplest to test from any terminal.",
 "MQTT — a lightweight publish/subscribe protocol over TCP/IP for bandwidth- and power-constrained devices.",
 "On IoTFlow both are first-class: a managed MQTT broker (port 1883) and a REST endpoint (/api/telemetry).",
 "Either way, every message is authenticated with your device token: Authorization: Bearer dev_..."],kicker="UPLINK",size=18)
tile_grid("MQTT in a Nutshell",[
 ("Publish / Subscribe","Clients publish messages to topics; subscribers receive them — no polling."),
 ("Broker","The central post office: dispatches every message from senders to the right receivers."),
 ("Topics","Routing keys such as devices/<id>/up — publish to a topic, subscribe to a topic."),
 ("Lightweight","Designed for constrained devices and unreliable networks; runs over TCP/IP or WebSockets.")],
 kicker="THE IOT PROTOCOL",cols=2,size=14)
img_slide("How MQTT Works","how-mqtt-works.png",
          "Publishers → broker → subscribers, routed by topic.",kicker="MQTT")
labs_cards(T); emit_labs(T)
content(f"Recap — {T['title']}",[
 "You registered a device on IoTFlow and secured its once-shown device token.",
 "You posted telemetry with cURL (HTTP), Python, and an ESP32 over MQTT.",
 "Telemetry is JSON metric key-value pairs; metric names bind to dashboard widgets.",
 "Readings appear in Latest Telemetry within seconds — auto-refresh every 5 s."],kicker="TOPIC RECAP",size=17)

# ================= TOPIC 3 =================
T=C.TOPICS[2]
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"ABILITIES & KNOWLEDGE: {T['tags']}",cols=2,size=14)
two_col("Reading Data from the Cloud",[
 ("REST API (pull)",0),("GET /api/device/state returns latest metrics & pins",1),
 ("Simple from any language — curl, Python requests",1),("Per-device telemetry history on the platform",1)],
 [("MQTT (push)",0),("Subscribe to your device's topics",1),
 ("Messages arrive the moment they are published",1),("Ideal for real-time apps and integrations",1)],
 kicker="A3 · INTEGRATE DATA SOURCES",lhead="HTTP GET",rhead="MQTT Subscribe")
tile_grid("Virtual Pins — Two-Way Control",[
 ("What they are","Named keys (V1, relay, pump) shared between dashboard widgets and device code — Blynk-style."),
 ("Control widgets","Button (momentary) · Switch (toggle) · Slider (value) · Terminal (text)."),
 ("Downlink path","Widget writes a pin → platform sends the command → device receives on devices/<id>/down (MQTT) or polls /api/device/state (HTTP)."),
 ("Device handler","onCommand(pin, value, text) maps the pin to a GPIO — relay, motor or LED responds instantly.")],
 kicker="REMOTE CONTROL",cols=2,size=14)
flow_h("Remote Control — the Downlink Path",[
 "User taps a Switch / moves a Slider on the dashboard",
 "The widget writes its virtual pin (e.g. V1 = 1)",
 "Platform delivers the command — MQTT downlink or HTTP poll",
 "Device onCommand() handler maps the pin to GPIO",
 "Relay / motor / LED responds — and events can fire automations"],kicker="END-TO-END",color=TEAL)
content("Alerts and Trigger Rules",[
 "Alert rules watch your telemetry: fire on a threshold (temperature > 30) or when a device goes offline.",
 "Active alerts are tracked per device and can be shown in an Alert-list dashboard widget.",
 "An alert can hand off to an n8n automation — notify, escalate, log or control a device back (Topic 4).",
 "Triggers are the bridge from monitoring to automation."],kicker="FROM DATA TO EVENTS",size=18)
labs_cards(T); emit_labs(T)
content(f"Recap — {T['title']}",[
 "You read device state over REST and streamed live messages with an MQTT subscription.",
 "You bound Control widgets to virtual pins and handled commands with onCommand().",
 "You flipped a relay/LED live from the dashboard — Blynk-style two-way control.",
 "Alert rules on thresholds or device-offline turn readings into triggers."],kicker="TOPIC RECAP",size=17)

brk("Lunch Break","1 hour")

# ================= TOPIC 4 =================
T=C.TOPICS[3]
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"ABILITIES & KNOWLEDGE: {T['tags']}",cols=2,size=14)
tile_grid("Dashboard Widgets",[
 ("Number card","The latest value of a metric, at a glance."),
 ("Gauge","A value against min/max bounds — great for %, temperature, tank level."),
 ("Line / Bar chart","Trends over time — the telemetry history plotted live."),
 ("LED indicator","Binary state — on/off, open/closed, healthy/alarm."),
 ("Map","Devices plotted by location — fleets, sites, sensors in the field."),
 ("Device status & alerts","Online/offline state and the list of active alerts.")],
 kicker="VISUALIZE (K5)",cols=2,size=14)
content("From Data to Business Insight",[
 "Raw readings become insight when you aggregate, trend and compare them — that is data analytics (K5).",
 "Dashboards answer 'what is happening now'; history charts answer 'what changed and when'.",
 "Automations answer 'what should happen next' — thresholds, schedules and AI recommendations.",
 "Review the data to produce insights of business value (A4): energy saved, downtime avoided, yield improved."],kicker="ANALYTICS",size=18)
tile_grid("What is n8n?",[
 ("Low-code automation","Drag-and-drop workflow editor — connect 400+ apps with zero code."),
 ("Nodes","Each node does one job: trigger, transform, notify, store, call AI."),
 ("Triggers & actions","A Webhook trigger starts the flow; action nodes do the work."),
 ("Workflows","Nodes wired together = a workflow. Activate it and it runs on every event.")],
 kicker="AUTOMATION ENGINE",cols=2,size=14)
flow_h("IoTFlow + n8n — Event-Driven Automation",[
 "Device event — telemetry, alert, online/offline or command",
 "IoTFlow Automation fires your n8n webhook URL",
 "n8n workflow runs — filter, branch, transform",
 "Actions — Email / Telegram / Sheets / AI / HTTP",
 "Optional: command sent back to the device"],kicker="HOW THEY CONNECT",color=VIOLET)
tile_grid("AI in Your IoT Workflows",[
 ("Summarise","An AI node turns a day of readings into a one-paragraph brief."),
 ("Detect anomalies","Ask the model whether the pattern looks abnormal before alerting."),
 ("Recommend","Generate an action: 'humidity rising — run the dehumidifier 2 h'."),
 ("Converse","Chat interfaces over your device data — ask questions in plain English.")],
 kicker="AI + IOT (K5)",cols=2,size=14)
img_slide("IoT Security — Why It Matters","iot-security.png",
          "Security for IoT adds complexity: the cyber and physical worlds converge.",kicker="K6 · CYBERSECURITY")
tile_grid("Threats to IoT Deployments",[
 ("Spoofing","An attacker manipulates device state anonymously or impersonates the originator (MitM)."),
 ("Tampering","Physical attacks — battery drain, sleep deprivation, RNG attacks."),
 ("Eavesdropping","Intercepting broadcasts to obtain information without authorisation."),
 ("Denial of Service","Jamming radio or cutting wires renders devices unable to communicate."),
 ("Elevation of privilege","A device is forced to do something else — a half-open valve tricked fully open."),
 ("Data breach","Weak credentials and unpatched firmware expose the whole network.")],
 kicker="K6 · CYBERSECURITY",cols=2,size=13)
content("Securing IoT Devices",[
 "Employ device discovery for complete visibility of what is connected.",
 "Apply network segmentation for a stronger defense.",
 "Adopt secure password practices — never keep default credentials; protect device tokens.",
 "Patch and update firmware whenever available.",
 "Actively monitor IoT devices at all times — alerts and security analytics."],kicker="K6 · CYBERSECURITY",size=18)
labs_cards(T); emit_labs(T)
content(f"Recap — {T['title']}",[
 "You built a live dashboard — cards, gauges, charts, LEDs and maps on a responsive grid.",
 "You wired IoTFlow events to an n8n webhook and built a no-code workflow.",
 "Your flow notified, logged to a sheet, called AI for insight and commanded the device back.",
 "You know the IoT threat landscape and the practices that secure a deployment."],kicker="TOPIC RECAP",size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("LO1 · IoT technologies","Explained IoT devices, sensors, actuators, triggers and wireless technologies."),
 ("LO2 · Post data to cloud","Registered a device and streamed telemetry over HTTP and MQTT."),
 ("LO3 · Control from cloud","Read data via REST/MQTT and controlled devices with virtual pins."),
 ("LO4 · Analytics & viz","Built dashboards and automated insight with n8n workflows and AI.")],
 kicker="LEARNING OUTCOMES",cols=2,size=15)
content("Keep Building",[
 f"Your IoTFlow account and projects remain yours — {C.PLATFORM_URL}.",
 f"Revisit the step-by-step tutorials at {C.TUTORIALS_URL}.",
 "Connect a real ESP32 or Raspberry Pi at home using the Integrate page snippets.",
 "Extend your n8n flow: add WhatsApp alerts, Google Sheets logging or an AI agent."],kicker="NEXT STEPS")
content("Recommended Courses",C.RECOMMENDED_COURSES,kicker="CONTINUE YOUR JOURNEY",size=17)
content("Support",[
 "If you have any enquiries during and after the class, you can contact us:",
 "Email: enquiry@tertiaryinfotech.com",
 "Tel: +65 6100 0613",
 f"Course page: {C.COURSE_URL}"],kicker="WE'RE HERE TO HELP")
content("Cert & TRAQOM Survey (Mandatory)",[
 "Complete the TRAQOM survey — scan the QR code on the LMS.",
 "Your e-certificate is issued via the LMS/TMS after you are assessed Competent.",
 "Portal: https://ai-lms-tms.tertiaryinfo.tech/"],kicker="BEFORE YOU LEAVE")
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance.",
 C.ASSESSMENT["note"]],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!","You are now ready to connect, visualise and automate the Internet of Things.","SEE YOU IN THE CLOUD",color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.FILE_STEM}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
